# core/exporter.py
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import os
import re

class UTF8PDF(FPDF):
    def __init__(self):
        super().__init__()
        try:
            # Try to use Arial Unicode MS if available (better Unicode support)
            self.add_font('Arial Unicode MS', '', 'arial-unicode-ms.ttf', uni=True)
            self.default_font = 'Arial Unicode MS'
        except Exception:
            # Fallback to Arial
            self.default_font = 'Arial'

    def sanitize_text(self, text):
        if not isinstance(text, str):
            text = str(text)
        # Specifically replace common non-latin-1 characters
        text = text.replace('→', '->').replace('•', '*').replace('…', '...').replace('−', '-')
        # Replace other problematic characters with similar ASCII ones
        text = text.encode('latin-1', 'replace').decode('latin-1')
        return text

    def safe_cell(self, w, h, txt='', border=0, ln=0, align='', fill=False):
        self.cell(w, h, self.sanitize_text(txt), border, ln, align, fill)

    def safe_multi_cell(self, w, h, txt='', border=0, align='', fill=False):
        self.multi_cell(w, h, self.sanitize_text(txt), border, align, fill)

def build_pdf(report_data: dict, report_title="Data Insights Report"):
    try:
        pdf = UTF8PDF()
    except Exception:
        pdf = FPDF()
    
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    # Title
    pdf.set_font("Arial", "B", 24)
    pdf.safe_cell(0, 20, report_title, ln=True, align="C")
    pdf.ln(10)
    
    # 1. Dataset Overview
    ov = report_data.get("overview", {})
    pdf.set_font("Arial", "B", 16)
    pdf.set_fill_color(240, 240, 240)
    pdf.safe_cell(0, 10, "1. Dataset Overview", ln=True, fill=True)
    pdf.ln(5)
    pdf.set_font("Arial", size=12)
    pdf.safe_multi_cell(0, 8, txt=ov.get("summary", ""))
    pdf.ln(5)
    
    if "dataset_name" in ov:
        pdf.set_font("Arial", "B", 12)
        pdf.safe_cell(0, 8, f"Dataset: {ov['dataset_name']}", ln=True)
    
    # 2. EDA Summary
    pdf.ln(5)
    pdf.set_font("Arial", "B", 16)
    pdf.safe_cell(0, 10, "2. Exploratory Data Analysis", ln=True, fill=True)
    pdf.ln(5)
    pdf.set_font("Arial", size=11)
    for k, v in report_data.get("eda", {}).items():
        pdf.set_font("Arial", "B", 11)
        pdf.safe_cell(0, 8, f"{k.title()}:", ln=True)
        pdf.set_font("Arial", size=11)
        pdf.safe_multi_cell(0, 8, "; ".join(v))
    
    # 3. Hypotheses
    pdf.ln(10)
    pdf.set_font("Arial", "B", 16)
    pdf.safe_cell(0, 10, "3. Hypothesis Testing Results", ln=True, fill=True)
    pdf.ln(5)
    for h in report_data.get("hypotheses", []):
        pdf.set_font("Arial", "B", 12)
        pdf.safe_cell(0, 8, f"- {h['title']}", ln=True)
        pdf.set_font("Arial", size=11)
        pdf.safe_multi_cell(0, 7, f"Test: {h['test']} | Result: {h['result']}")
        pdf.safe_multi_cell(0, 7, f"Interpretation: {h['interpretation']}")
        pdf.ln(3)

    # 4. AI Recommendations
    pdf.ln(5)
    pdf.set_font("Arial", "B", 16)
    pdf.safe_cell(0, 10, "4. Suggested Next Steps", ln=True, fill=True)
    pdf.ln(5)
    pdf.set_font("Arial", size=11)
    for s in report_data.get("suggestions", []):
        pdf.safe_multi_cell(0, 8, f"* {s}")

    # 5. Visualizations (Key Page)
    if report_data.get("figures"):
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.safe_cell(0, 10, "5. Visual Data Evidence", ln=True, fill=True)
        pdf.ln(10)
        
        for i, fig_path in enumerate(report_data["figures"]):
            if os.path.exists(fig_path):
                # Add 2 images per page
                if i > 0 and i % 2 == 0:
                    pdf.add_page()
                try:
                    pdf.image(fig_path, x=15, w=180)
                    pdf.ln(5)
                except Exception as e:
                    pdf.safe_cell(0, 10, f"[Error adding figure {i}]", ln=True)

    # Save
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf.output(tmp.name)
    return tmp.name


def build_ppt(report_data: dict, report_title="Data Insights Report"):
    prs = Presentation()

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = report_title
    subtitle = slide.placeholders[1]
    ov = report_data.get("overview", {})
    subtitle.text = f"Dataset: {ov.get('dataset_name', 'Unknown')}\nSize: {ov.get('shape', 'N/A')}\nCompiled via InSightGenie"

    # 2. Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Overview"
    body = slide.placeholders[1]
    body.text = ov.get("summary", "No summary available.")

    # 3. EDA Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exploratory Insights"
    body = slide.placeholders[1]
    eda_lines = []
    for k, v in report_data.get("eda", {}).items():
        eda_lines.append(f"{k.title()}: " + "; ".join(v))
    body.text = "\n".join(eda_lines)

    # 4. Hypothesis Testing
    hypos = report_data.get("hypotheses", [])
    if hypos:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Statistical Hypotheses"
        body = slide.placeholders[1]
        body.text = "\n".join([f"{h['title']}: {h['result']}" for h in hypos[:8]])

    # 5. Visualizations
    for i, fig_path in enumerate(report_data.get("figures", [])):
        if os.path.exists(fig_path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            slide.shapes.title.text = f"Visualization Proof {i+1}"
            left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(5)
            slide.shapes.add_picture(fig_path, left, top, width, height)

    # 6. AI Suggestions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Suggested Analyses"
    body = slide.placeholders[1]
    body.text = "\n".join(report_data.get("suggestions", []))

    # Save
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name
